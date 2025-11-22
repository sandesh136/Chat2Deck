from .gemini_client import generate_slides
from .ppt_generator import create_pptx_from_code
import logging

def query_to_pptx(query: str, num_slides: int):
    example_code = '''
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "🚀 Amazing AI-Generated Presentation"
slide.placeholders[1].text = "Generated using python-pptx ✨\nEmpower your storytelling with data & design"

# Content Slide with bullet points and emojis
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Key Features ✅"
content = slide.placeholders[1].text_frame
content.clear()
points = [
    "💡 AI-powered slide generation",
    "📊 Interactive charts with crisp visuals",
    "🎨 Professional design with shapes & colors",
    "🔗 Supports live data and dynamic updates"
]
for point in points:
    p = content.add_paragraph()
    p.text = point
    p.level = 0
    p.font.size = Pt(18)

# Chart Slide
chart_slide_layout = prs.slide_layouts[5]  # Title Only layout
slide = prs.slides.add_slide(chart_slide_layout)
slide.shapes.title.text = "Sales Data Overview 📈"

chart_data = CategoryChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data.add_series('Product A', (19.5, 21.0, 16.8, 24.0))
chart_data.add_series('Product B', (14.0, 18.5, 23.0, 22.5))

x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

# Customize chart series colors
series_colors = [RGBColor(0x4F, 0x81, 0xBD), RGBColor(0xC0, 0x50, 0x4D)]
for i, series in enumerate(chart.series):
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = series_colors[i]

# Add a Fancy Shape Callout
shape = slide.shapes.add_shape(MSO_SHAPE.CLOUD_CALLOUT, Inches(0.2), Inches(1), Inches(3), Inches(1))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)
shape.text_frame.text = "Dynamic & Engaging! 🚀"
for paragraph in shape.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Add another slide with a diagram made of connected shapes
diagram_slide = prs.slides.add_slide(prs.slide_layouts[5])
diagram_slide.shapes.title.text = "Process Flow Diagram 🔄"

left = Inches(1)
top = Inches(2)
width = Inches(2)
height = Inches(1)

# Shapes for process steps
step1 = diagram_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
step1.text = "Start"
step1.fill.solid()
step1.fill.fore_color.rgb = RGBColor(0x67, 0xA9, 0x41)  # green
step1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
step1.text_frame.paragraphs[0].font.bold = True

step2 = diagram_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(3), top, width, height)
step2.text = "Process"
step2.fill.solid()
step2.fill.fore_color.rgb = RGBColor(0xF7, 0xB5, 0x4D)  # orange
step2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
step2.text_frame.paragraphs[0].font.bold = True

step3 = diagram_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(6), top, width, height)
step3.text = "End"
step3.fill.solid()
step3.fill.fore_color.rgb = RGBColor(0x4A, 0x90, 0xE2)  # blue
step3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
step3.text_frame.paragraphs[0].font.bold = True

# Connectors (arrows)
line1 = slide.shapes.add_connector(
    MSO_CONNECTOR_TYPE.ELBOW, left + width, top + height/2, left + Inches(3), top + height/2
)
line2 = slide.shapes.add_connector(
    MSO_CONNECTOR_TYPE.ELBOW, left + Inches(5), top + height/2, left + Inches(6), top + height/2
)

prs.save("generated_presentation.pptx")

'''

    initial_prompt = (
        f"You are an expert assistant that generates Python code using the 'python-pptx' library "
        f"to create a professional, interactive PowerPoint presentation titled '{query}' "
        f"with approximately {num_slides} slides. "
        "The Python code must be complete, runnable, and self-contained. "
        "Include high-quality charts with labels and legends, text formatting, bullet points with emojis, and polished styling. "
        "Below is an example of a well-structured, runnable python-pptx script that your output should emulate in style and completeness. "
        "The code must save the presentation to 'generated_presentation.pptx' exactly as shown in the example:\n"
        f"{example_code}\n"
        "Return only the complete Python code without any explanations or additional formatting."
    )

    max_retries = 3
    last_error = None
    python_code = None

    for attempt in range(max_retries):
        logging.info(f"Presentation generation attempt {attempt + 1} of {max_retries}")
        try:
            prompt = initial_prompt if attempt == 0 else (
                f"The previous Python code you generated failed with an error. Please fix it.\n\n"
                f"Original Request: Create a presentation about '{query}'.\n\n"
                f"--- BROKEN CODE ---\n{python_code}\n\n"
                f"--- ERROR MESSAGE ---\n{last_error}\n\n"
                f"Provide the complete, corrected Python code that fixes the error. Do not include explanations."
            )
            python_code = generate_slides(prompt)
            return create_pptx_from_code(python_code)
        except Exception as e:
            last_error = e
            logging.warning(f"Attempt {attempt + 1} failed with error: {e}")

    raise RuntimeError(f"Failed to generate a valid presentation after {max_retries} attempts. Last error: {last_error}")
